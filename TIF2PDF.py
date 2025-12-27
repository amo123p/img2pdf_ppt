#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
图片转PDF/PPT转换器 - GUI版本

功能描述：
    将任意图片文件合并转换为PDF或PPT文档，
    支持多种图片格式，提供直观的图形界面操作。

依赖安装：
    pip install pillow pymupdf python-pptx

创建时间：2025-12-27
"""

import os
import sys
import glob
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Optional, Tuple
import io


# 支持的图片格式
SUPPORTED_FORMATS = {
    '.jpg', '.jpeg', '.png', '.bmp', '.gif',
    '.tif', '.tiff', '.webp', '.ico', '.ppm'
}


class ImageConverterGUI:
    """图片转换器GUI主类"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("图片转PDF/PPT转换器")
        self.root.geometry("900x700")
        self.root.minsize(800, 600)
        
        # 数据存储
        self.image_files = []
        self.selected_format = tk.StringVar(value="PDF")
        
        # 设置样式
        self.setup_styles()
        
        # 创建界面
        self.create_widgets()
        
        # 绑定事件
        self.bind_events()
    
    def setup_styles(self):
        """设置界面样式"""
        style = ttk.Style()
        style.configure('TButton', font=('微软雅黑', 10))
        style.configure('TLabel', font=('微软雅黑', 10))
        style.configure('Header.TLabel', font=('微软雅黑', 12, 'bold'))
        style.configure('TRadiobutton', font=('微软雅黑', 10))
    
    def create_widgets(self):
        """创建所有界面组件"""
        # 主容器
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # ==================== 顶部：文件夹选择区域 ====================
        top_frame = ttk.LabelFrame(main_frame, text=" 1. 选择图片文件夹 ", padding="10")
        top_frame.pack(fill=tk.X, pady=(0, 10))
        
        # 文件夹路径输入框
        path_frame = ttk.Frame(top_frame)
        path_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(path_frame, text="文件夹路径:").pack(side=tk.LEFT)
        
        self.path_entry = ttk.Entry(path_frame, width=70)
        self.path_entry.pack(side=tk.LEFT, padx=(10, 5), fill=tk.X, expand=True)
        
        self.browse_btn = ttk.Button(path_frame, text="浏览...", command=self.browse_folder)
        self.browse_btn.pack(side=tk.LEFT, padx=(5, 0))
        
        # 图片格式过滤器
        filter_frame = ttk.Frame(top_frame)
        filter_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(filter_frame, text="图片格式过滤:").pack(side=tk.LEFT)
        
        self.format_var = tk.StringVar(value="所有图片")
        format_combo = ttk.Combobox(filter_frame, textvariable=self.format_var,
                                     values=["所有图片", "仅PNG/JPG", "仅TIF", "仅BMP"],
                                     width=20, state="readonly")
        format_combo.pack(side=tk.LEFT, padx=(10, 0))
        
        self.rescan_btn = ttk.Button(filter_frame, text="刷新列表", command=self.scan_folder)
        self.rescan_btn.pack(side=tk.LEFT, padx=(10, 0))
        
        # ==================== 中部：图片列表区域 ====================
        mid_frame = ttk.LabelFrame(main_frame, text=" 2. 图片列表（可拖拽调整顺序） ", padding="10")
        mid_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        # 列表控制按钮
        list_btn_frame = ttk.Frame(mid_frame)
        list_btn_frame.pack(fill=tk.X, pady=(0, 5))
        
        self.select_all_btn = ttk.Button(list_btn_frame, text="全选", command=self.select_all)
        self.select_all_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        self.deselect_all_btn = ttk.Button(list_btn_frame, text="取消全选", command=self.deselect_all)
        self.deselect_all_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        self.reverse_btn = ttk.Button(list_btn_frame, text="反选", command=self.reverse_selection)
        self.reverse_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        self.move_up_btn = ttk.Button(list_btn_frame, text="上移", command=self.move_up)
        self.move_up_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        self.move_down_btn = ttk.Button(list_btn_frame, text="下移", command=self.move_down)
        self.move_down_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        self.remove_btn = ttk.Button(list_btn_frame, text="移除选中", command=self.remove_selected)
        self.remove_btn.pack(side=tk.LEFT, padx=(0, 5))
        
        ttk.Label(list_btn_frame, text="  图片数量:").pack(side=tk.LEFT, padx=(20, 0))
        self.count_label = ttk.Label(list_btn_frame, text="0", foreground="blue")
        self.count_label.pack(side=tk.LEFT)
        
        # 图片列表（带滚动条）
        list_frame = ttk.Frame(mid_frame)
        list_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建带复选框的列表
        self.listbox = tk.Listbox(list_frame, selectmode=tk.EXTENDED,
                                   font=('Consolas', 10), height=10)
        self.listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 滚动条
        scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.listbox.config(yscrollcommand=scrollbar.set)
        
        # ==================== 底部：输出设置区域 ====================
        bottom_frame = ttk.LabelFrame(main_frame, text=" 3. 输出设置 ", padding="10")
        bottom_frame.pack(fill=tk.X, pady=(0, 10))
        
        # 输出格式选择
        format_frame = ttk.Frame(bottom_frame)
        format_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(format_frame, text="输出格式:").pack(side=tk.LEFT)
        
        ttk.Radiobutton(format_frame, text="PDF文档", variable=self.selected_format,
                        value="PDF", command=self.on_format_change).pack(side=tk.LEFT, padx=(10, 5))
        
        ttk.Radiobutton(format_frame, text="PPT演示文稿", variable=self.selected_format,
                        value="PPT", command=self.on_format_change).pack(side=tk.LEFT, padx=(5, 0))
        
        # 输出路径
        output_frame = ttk.Frame(bottom_frame)
        output_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(output_frame, text="输出路径:").pack(side=tk.LEFT)
        
        self.output_entry = ttk.Entry(output_frame, width=70)
        self.output_entry.pack(side=tk.LEFT, padx=(10, 5), fill=tk.X, expand=True)
        
        self.output_browse_btn = ttk.Button(output_frame, text="浏览...", command=self.browse_output)
        self.output_browse_btn.pack(side=tk.LEFT, padx=(5, 0))
        
        # 选项设置
        options_frame = ttk.Frame(bottom_frame)
        options_frame.pack(fill=tk.X, pady=5)
        
        # PDF选项
        self.pdf_quality_var = tk.StringVar(value="无损")
        self.pdf_quality_combo = ttk.Combobox(options_frame, textvariable=self.pdf_quality_var,
                                               values=["无损（推荐）", "中等质量", "高压缩"],
                                               width=20, state="readonly")
        self.pdf_quality_combo.pack(side=tk.LEFT, padx=(10, 0))
        
        # PPT选项
        self.ppt_layout_var = tk.StringVar(value="空白页")
        self.ppt_layout_combo = ttk.Combobox(options_frame, textvariable=self.ppt_layout_var,
                                              values=["空白页", "图片居中", "全屏填充"],
                                              width=20, state="readonly")
        self.ppt_layout_combo.pack(side=tk.LEFT, padx=(10, 0))
        self.ppt_layout_combo.pack_forget()  # 初始隐藏
        
        # ==================== 转换按钮和进度区域 ====================
        action_frame = ttk.Frame(main_frame)
        action_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.convert_btn = ttk.Button(action_frame, text=" 开始转换 ",
                                       command=self.start_conversion, width=20)
        self.convert_btn.pack(side=tk.TOP)
        
        # 进度条
        progress_frame = ttk.Frame(main_frame)
        progress_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var,
                                            maximum=100, mode='determinate')
        self.progress_bar.pack(fill=tk.X, pady=(0, 5))
        
        self.status_label = ttk.Label(progress_frame, text="就绪", foreground="green")
        self.status_label.pack(side=tk.LEFT)
        
        # ==================== 日志输出区域 ====================
        log_frame = ttk.LabelFrame(main_frame, text=" 日志输出 ", padding="5")
        log_frame.pack(fill=tk.BOTH, expand=True)
        
        self.log_text = scrolledtext.ScrolledText(log_frame, height=8, font=('Consolas', 9))
        self.log_text.pack(fill=tk.BOTH, expand=True)
    
    def bind_events(self):
        """绑定事件处理"""
        self.listbox.bind('<Double-Button-1>', self.on_listbox_double_click)
        self.listbox.bind('<<ListboxSelect>>', self.on_listbox_select)
    
    def browse_folder(self):
        """浏览选择文件夹"""
        folder = filedialog.askdirectory(title="选择包含图片的文件夹")
        if folder:
            self.path_entry.delete(0, tk.END)
            self.path_entry.insert(0, folder)
            self.scan_folder()
    
    def browse_output(self):
        """浏览选择输出路径"""
        format_type = self.selected_format.get()
        
        if format_type == "PDF":
            filetypes = [("PDF文件", "*.pdf")]
            default_ext = ".pdf"
        else:
            filetypes = [("PowerPoint演示文稿", "*.pptx")]
            default_ext = ".pptx"
        
        initial_dir = self.path_entry.get() or os.path.expanduser("~")
        output_file = filedialog.asksaveasfilename(
            title="选择输出路径",
            filetypes=filetypes,
            defaultextension=default_ext,
            initialdir=initial_dir
        )
        
        if output_file:
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, output_file)
    
    def scan_folder(self):
        """扫描文件夹中的图片"""
        folder = self.path_entry.get().strip()
        if not folder or not os.path.isdir(folder):
            self.log("错误：请先选择有效的文件夹路径")
            return
        
        format_filter = self.format_var.get()
        self.image_files = []
        
        # 扫描图片文件
        for ext in SUPPORTED_FORMATS:
            if format_filter == "仅PNG/JPG" and ext not in ['.jpg', '.jpeg', '.png']:
                continue
            elif format_filter == "仅TIF" and ext not in ['.tif', '.tiff']:
                continue
            elif format_filter == "仅BMP" and ext not in ['.bmp']:
                continue
            
            pattern = os.path.join(folder, f'*{ext}')
            self.image_files.extend(glob.glob(pattern))
        
        # 去重并排序
        self.image_files = sorted(list(set(self.image_files)))
        
        # 更新列表显示
        self.update_listbox()
        
        self.log(f"扫描完成：找到 {len(self.image_files)} 个图片文件")
        
        # 自动设置输出路径
        if self.image_files:
            first_img = self.image_files[0]
            folder_name = os.path.splitext(os.path.basename(first_img))[0]
            output_dir = os.path.dirname(first_img)
            
            format_ext = ".pdf" if self.selected_format.get() == "PDF" else ".pptx"
            output_path = os.path.join(output_dir, f"{folder_name}_合并{format_ext}")
            
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, output_path)
    
    def update_listbox(self):
        """更新图片列表显示"""
        self.listbox.delete(0, tk.END)
        
        for i, filepath in enumerate(self.image_files):
            filename = os.path.basename(filepath)
            file_size = os.path.getsize(filepath) / (1024 * 1024)  # MB
            
            # 尝试获取图片尺寸
            try:
                with Image.open(filepath) as img:
                    width, height = img.size
                    size_info = f"  [{width}x{height}]"
            except:
                size_info = ""
            
            display_text = f"{i+1:3d}. {filename} ({file_size:.2f} MB){size_info}"
            self.listbox.insert(tk.END, display_text)
        
        self.count_label.config(text=str(len(self.image_files)))
    
    def on_listbox_double_click(self, event):
        """双击列表项"""
        selection = self.listbox.curselection()
        if selection:
            index = selection[0]
            if index < len(self.image_files):
                filepath = self.image_files[index]
                try:
                    os.startfile(filepath)
                except:
                    self.log(f"无法打开文件: {filepath}")
    
    def on_listbox_select(self, event):
        """列表选择事件"""
        pass
    
    def select_all(self):
        """全选"""
        self.listbox.selection_set(0, tk.END)
    
    def deselect_all(self):
        """取消全选"""
        self.listbox.selection_clear(0, tk.END)
    
    def reverse_selection(self):
        """反选"""
        all_items = list(range(self.listbox.size()))
        selected = set(self.listbox.curselection())
        for i in all_items:
            if i not in selected:
                self.listbox.selection_set(i)
    
    def move_up(self):
        """上移选中项"""
        selected = list(self.listbox.curselection())
        if not selected or selected[0] == 0:
            return
        
        for i in selected:
            if i > 0:
                self.image_files[i], self.image_files[i-1] = self.image_files[i-1], self.image_files[i]
        
        self.update_listbox()
        for i in selected:
            self.listbox.selection_set(i - 1 if i > 0 else i)
    
    def move_down(self):
        """下移选中项"""
        selected = list(self.listbox.curselection())
        if not selected or selected[-1] == len(self.image_files) - 1:
            return
        
        for i in reversed(selected):
            if i < len(self.image_files) - 1:
                self.image_files[i], self.image_files[i+1] = self.image_files[i+1], self.image_files[i]
        
        self.update_listbox()
        for i in selected:
            self.listbox.selection_set(i + 1 if i < len(self.image_files) - 1 else i)
    
    def remove_selected(self):
        """移除选中的项目"""
        selected = list(self.listbox.curselection())
        if not selected:
            return
        
        # 从后往前删除
        for i in reversed(selected):
            if i < len(self.image_files):
                del self.image_files[i]
        
        self.update_listbox()
        self.log(f"已移除 {len(selected)} 个项目")
    
    def on_format_change(self):
        """输出格式改变时的处理"""
        format_type = self.selected_format.get()
        
        if format_type == "PDF":
            self.pdf_quality_combo.pack(side=tk.LEFT, padx=(10, 0))
            self.ppt_layout_combo.pack_forget()
            
            # 更新输出路径扩展名
            current_output = self.output_entry.get()
            if current_output:
                new_output = os.path.splitext(current_output)[0] + ".pdf"
                self.output_entry.delete(0, tk.END)
                self.output_entry.insert(0, new_output)
        else:
            self.pdf_quality_combo.pack_forget()
            self.ppt_layout_combo.pack(side=tk.LEFT, padx=(10, 0))
            
            # 更新输出路径扩展名
            current_output = self.output_entry.get()
            if current_output:
                new_output = os.path.splitext(current_output)[0] + ".pptx"
                self.output_entry.delete(0, tk.END)
                self.output_entry.insert(0, new_output)
    
    def start_conversion(self):
        """开始转换（在独立线程中执行）"""
        # 验证输入
        if not self.image_files:
            messagebox.showwarning("警告", "没有可转换的图片文件！")
            return
        
        output_file = self.output_entry.get().strip()
        if not output_file:
            messagebox.showwarning("警告", "请设置输出路径！")
            return
        
        # 禁用按钮防止重复点击
        self.convert_btn.config(state=tk.DISABLED)
        self.progress_var.set(0)
        
        # 在新线程中执行转换
        thread = threading.Thread(target=self.convert_images, args=(output_file,))
        thread.daemon = True
        thread.start()
    
    def convert_images(self, output_file):
        """执行转换（在线程中运行）"""
        format_type = self.selected_format.get()
        
        try:
            if format_type == "PDF":
                self.convert_to_pdf(output_file)
            else:
                self.convert_to_ppt(output_file)
            
            self.log(f"转换完成！输出文件: {output_file}")
            self.status_label.config(text="转换成功！", foreground="green")
            messagebox.showinfo("完成", f"转换成功！\n输出文件: {output_file}")
            
        except Exception as e:
            self.log(f"转换失败: {str(e)}")
            self.status_label.config(text="转换失败", foreground="red")
            messagebox.showerror("错误", f"转换失败！\n{str(e)}")
        
        finally:
            # 恢复按钮状态
            self.root.after(0, lambda: self.convert_btn.config(state=tk.NORMAL))
    
    def convert_to_pdf(self, output_file):
        """转换为PDF"""
        quality_mode = self.pdf_quality_var.get()
        self.log(f"开始转换为PDF（{quality_mode}）...")

        # 根据压缩模式设置参数
        if quality_mode == "无损（推荐）":
            compression_type = "lossless"
            quality = None
        elif quality_mode == "中等质量":
            compression_type = "compress"
            quality = 60
        else:  # 高压缩（小文件）
            compression_type = "compress"
            quality = 30

        # 创建PDF文档
        pdf_document = fitz.open()

        total_files = len(self.image_files)
        total_size = 0

        for i, filepath in enumerate(self.image_files):
            filename = os.path.basename(filepath)
            file_size = os.path.getsize(filepath)
            total_size += file_size

            self.log(f"处理 ({i+1}/{total_files}): {filename} [{quality_mode}]")
            self.status_label.config(text=f"正在处理: {filename}", foreground="blue")

            try:
                with Image.open(filepath) as img:
                    width, height = img.size

                    if compression_type == "lossless":
                        # 无损模式：直接嵌入原始图片
                        page = pdf_document.new_page(width=width, height=height)
                        page.insert_image(page.rect, filename=filepath, overlay=False)
                    else:
                        # 压缩模式：先压缩再嵌入
                        # 转换为RGB模式（PDF不支持RGBA）
                        if img.mode == 'RGBA':
                            img = img.convert('RGB')
                        elif img.mode == 'CMYK':
                            img = img.convert('RGB')

                        # 保存为压缩的JPEG数据
                        img_buffer = io.BytesIO()
                        img.save(img_buffer, format='JPEG', quality=quality, optimize=True)
                        img_data = img_buffer.getvalue()

                        # 创建PDF页面并插入压缩后的图片
                        page = pdf_document.new_page(width=width, height=height)
                        page.insert_image(page.rect, stream=img_data, overlay=False)

            except Exception as e:
                self.log(f"  警告：无法处理图片: {e}")
                continue

            # 更新进度
            progress = ((i + 1) / total_files) * 100
            self.root.after(0, lambda p=progress: self.progress_var.set(p))

        # 保存PDF
        self.log("正在保存PDF文件...")
        pdf_document.save(output_file)
        pdf_document.close()

        output_size = os.path.getsize(output_file)
        compression_ratio = (1 - output_size / total_size) * 100 if total_size > 0 else 0

        self.log(f"PDF生成完成！")
        self.log(f"  原始总大小: {total_size / (1024*1024):.2f} MB")
        self.log(f"  PDF大小: {output_size / (1024*1024):.2f} MB")
        self.log(f"  压缩比: {compression_ratio:.1f}%")
    
    def convert_to_ppt(self, output_file):
        """转换为PPT"""
        self.log("开始转换为PPT...")
        
        # 创建演示文稿
        prs = Presentation()
        
        total_files = len(self.image_files)
        ppt_layout = self.ppt_layout_var.get()
        
        for i, filepath in enumerate(self.image_files):
            filename = os.path.basename(filepath)
            file_size = os.path.getsize(filepath) / (1024 * 1024)
            
            self.log(f"处理 ({i+1}/{total_files}): {filename}")
            self.status_label.config(text=f"正在处理: {filename}", foreground="blue")
            
            # 获取图片尺寸
            try:
                with Image.open(filepath) as img:
                    original_width, original_height = img.size
            except Exception as e:
                self.log(f"  警告：无法读取图片尺寸: {e}")
                continue
            
            # 添加新幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 根据选择的布局处理图片
            if ppt_layout == "全屏填充":
                # 全屏填充，保持比例
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                ratio = min(slide_width / original_width, slide_height / original_height)
                img_width = int(original_width * ratio)
                img_height = int(original_height * ratio)
                
                left = int((slide_width - img_width) / 2)
                top = int((slide_height - img_height) / 2)
                
            elif ppt_layout == "图片居中":
                # 居中显示，保持原始大小
                img_width = int(original_width)
                img_height = int(original_height)
                
                left = int((prs.slide_width - img_width) / 2)
                top = int((prs.slide_height - img_height) / 2)
                
            else:  # 空白页
                # 直接使用原始尺寸
                img_width = int(original_width)
                img_height = int(original_height)
                left = Inches(0.5)
                top = Inches(0.5)
            
            # 插入图片
            try:
                slide.shapes.add_picture(filepath, left, top, width=img_width, height=img_height)
            except Exception as e:
                self.log(f"  警告：插入图片失败: {e}")
                continue
            
            # 更新进度
            progress = ((i + 1) / total_files) * 100
            self.root.after(0, lambda p=progress: self.progress_var.set(p))
        
        # 保存PPT
        self.log("正在保存PPT文件...")
        prs.save(output_file)
        
        output_size = os.path.getsize(output_file)
        self.log(f"PPT生成完成！")
        self.log(f"  包含 {total_files} 张幻灯片")
        self.log(f"  PPT大小: {output_size / (1024*1024):.2f} MB")
    
    def log(self, message):
        """添加日志"""
        self.root.after(0, lambda: self.log_text.insert(tk.END, message + "\n"))
        self.root.after(0, lambda: self.log_text.see(tk.END))


def main():
    """主函数"""
    root = tk.Tk()
    
    # 设置窗口图标（如果有的话）
    try:
        root.iconbitmap(default=None)
    except:
        pass
    
    app = ImageConverterGUI(root)
    root.mainloop()


if __name__ == '__main__':
    main()
